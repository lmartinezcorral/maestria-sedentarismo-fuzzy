#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GENERADOR DE PRESENTACIÃ“N POWERPOINT
Sistema de Inferencia Difusa para EvaluaciÃ³n de Sedentarismo
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# ============================================================================
# CONFIGURACIÃ“N
# ============================================================================

BASE_DIR = Path(__file__).parent
OUTPUT_FILE = BASE_DIR / 'PRESENTACION_SISTEMA_DIFUSO_SEDENTARISMO.pptx'
PLOTS_DIR = BASE_DIR / 'analisis_u' / 'fuzzy' / 'plots'

# Colores del tema
COLOR_TITULO = RGBColor(0, 51, 102)  # Azul oscuro
COLOR_SUBTITULO = RGBColor(51, 102, 153)  # Azul medio
COLOR_TEXTO = RGBColor(64, 64, 64)  # Gris oscuro
COLOR_ACENTO = RGBColor(0, 176, 80)  # Verde para highlights

print("=" * 80)
print("GENERADOR DE PRESENTACIÃ“N POWERPOINT")
print("=" * 80)
print()

# ============================================================================
# CREAR PRESENTACIÃ“N
# ============================================================================

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================================================
# FUNCIONES AUXILIARES
# ============================================================================


def add_title_slide(title, subtitle=""):
    """Agrega slide de tÃ­tulo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # TÃ­tulo
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_TITULO
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # SubtÃ­tulo
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = COLOR_SUBTITULO
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(title):
    """Agrega slide con tÃ­tulo y retorna el slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # TÃ­tulo
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    # LÃ­nea separadora
    line = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(1), Inches(9), Inches(0)
    )
    line.line.color.rgb = COLOR_ACENTO
    line.line.width = Pt(3)

    return slide


def add_bullet_text(slide, left, top, width, height, text_lines, font_size=18):
    """Agrega cuadro de texto con bullets"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(text_lines):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXTO
        p.level = 0 if not line.startswith("  ") else 1

    return text_box


def add_image(slide, image_path, left, top, width=None, height=None):
    """Agrega imagen al slide"""
    if not image_path.exists():
        print(f"   [ADVERTENCIA] Imagen no encontrada: {image_path}")
        return None

    return slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def add_table(slide, left, top, rows, cols, data):
    """Agrega tabla al slide"""
    table_shape = slide.shapes.add_table(
        rows, cols, left, top, Inches(8), Inches(1.5))
    table = table_shape.table

    # Llenar tabla
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.rows[i].cells[j]
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(14)

            # Encabezado en negrita
            if i == 0:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 220, 240)

    return table_shape


# ============================================================================
# SLIDE 1: TÃTULO
# ============================================================================

print("[1/15] Generando slide de tÃ­tulo...")
add_title_slide(
    "Sistema de Inferencia Difusa para EvaluaciÃ³n del Comportamiento Sedentario",
    "Modelo Mamdani con ValidaciÃ³n vs. Clustering No Supervisado"
)

# Agregar autor y fecha
slide = prs.slides[-1]
author_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
author_frame = author_box.text_frame
author_frame.text = "Luis Ãngel MartÃ­nez\nMaestrÃ­a en Ciencias, Semestre 3\n18 de octubre de 2025"
author_frame.paragraphs[0].font.size = Pt(18)
author_frame.paragraphs[0].font.color.rgb = COLOR_TEXTO
author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 2: OBJETIVOS
# ============================================================================

print("[2/15] Generando slide de objetivos...")
slide = add_content_slide("Objetivos")

add_bullet_text(
    slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5),
    [
        "OBJETIVO PRINCIPAL:",
        "  Desarrollar y validar un sistema de inferencia difusa tipo Mamdani para clasificar el sedentarismo semanal a partir de biomarcadores de wearables (Apple Watch)",
        "",
        "OBJETIVOS ESPECÃFICOS:",
        "  â€¢ Preprocesar datos con imputaciÃ³n jerÃ¡rquica sin leak temporal",
        "  â€¢ Crear variables derivadas normalizadas (Actividad_relativa, SuperÃ¡vit_calÃ³rico_basal)",
        "  â€¢ Descubrir estructura latente con clustering K-means (K=2)",
        "  â€¢ DiseÃ±ar funciones de membresÃ­a triangulares por percentiles",
        "  â€¢ Definir reglas difusas clÃ­nicamente interpretables",
        "  â€¢ Validar sistema fuzzy vs. clustering (bÃºsqueda de umbral Ï„ Ã³ptimo)"
    ],
    font_size=20
)

# ============================================================================
# SLIDE 3: DATOS Y COHORTE
# ============================================================================

print("[3/15] Generando slide de cohorte...")
slide = add_content_slide("Datos y Cohorte")

add_bullet_text(
    slide, Inches(0.8), Inches(1.5), Inches(4), Inches(2),
    [
        "â€¢ Cohorte: 10 adultos (5M/5H)",
        "â€¢ Seguimiento: Multianual",
        "â€¢ Unidad de anÃ¡lisis: 1,337 semanas vÃ¡lidas",
        "â€¢ Fuente: Apple Watch (datos diarios)",
        "â€¢ Variables base: Actividad, FC, HRV, gasto calÃ³rico"
    ],
    font_size=18
)

# Tabla de caracterÃ­sticas
table_data = [
    ["Usuario", "Sexo", "Edad", "Peso (kg)", "TMB (kcal/d)", "Semanas"],
    ["u1 (ale)", "M", "34", "68", "1411", "149"],
    ["u6 (fidel)", "H", "34", "100", "1958", "278"],
    ["u9 (lmartinez)", "H", "32", "124", "2241", "298"],
    ["u10 (vane)", "M", "28", "58", "1304", "131"],
    ["...", "...", "...", "...", "...", "..."]
]

add_table(slide, Inches(5.2), Inches(1.8), 6, 6, table_data)

# ============================================================================
# SLIDE 4: VARIABLES CLAVE
# ============================================================================

print("[4/15] Generando slide de variables clave...")
slide = add_content_slide("Variables Derivadas Clave")

add_bullet_text(
    slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5),
    [
        "1. ACTIVIDAD_RELATIVA (normalizada por exposiciÃ³n):",
        "   = min_movimiento / (60 Ã— hrs_monitoreadas)",
        "   Rationale: Corrige por tiempo de uso del reloj",
        "",
        "2. SUPERÃVIT_CALÃ“RICO_BASAL (ajustado por antropometrÃ­a):",
        "   = (Gasto_activo Ã— 100) / TMB",
        "   Rationale: 400 kcal â‰  impacto equivalente (depende de metabolismo basal)",
        "",
        "3. HRV_SDNN (variabilidad cardÃ­aca, marcador vagal):",
        "   Alta HRV (>60 ms) = tono vagal saludable",
        "   Baja HRV (<40 ms) = desacondicionamiento, estrÃ©s",
        "",
        "4. DELTA_CARDIACO (respuesta al esfuerzo):",
        "   = FC_caminata - FC_reposo"
    ],
    font_size=17
)

# ============================================================================
# SLIDE 5: PIPELINE METODOLÃ“GICO
# ============================================================================

print("[5/15] Generando slide de pipeline...")
slide = add_content_slide("Pipeline MetodolÃ³gico (5 Fases)")

add_bullet_text(
    slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5),
    [
        "FASE 1: Preprocesamiento Diario + ImputaciÃ³n JerÃ¡rquica",
        "  â€¢ Gates: Hard no-wear, Soft low-activity, Normal",
        "  â€¢ Rolling mediana (solo pasado, sin leak temporal)",
        "",
        "FASE 2: CreaciÃ³n de Variables Derivadas",
        "  â€¢ Actividad_relativa, SuperÃ¡vit_calÃ³rico_basal (reemplazan originales)",
        "",
        "FASE 3: AgregaciÃ³n Semanal Robusta",
        "  â€¢ 8 features: p50 (mediana) + IQR por variable clave",
        "",
        "FASE 4: Clustering No Supervisado (Verdad Operativa)",
        "  â€¢ K-Means con K-sweep (K=2..6), K=2 Ã³ptimo (Silhouette=0.232)",
        "",
        "FASE 5: Sistema de Inferencia Difusa + ValidaciÃ³n",
        "  â€¢ 4 entradas Ã— 3 etiquetas, 5 reglas Mamdani, Ï„=0.30"
    ],
    font_size=17
)

# ============================================================================
# SLIDE 6: CLUSTERING K=2
# ============================================================================

print("[6/15] Generando slide de clustering...")
slide = add_content_slide("Clustering K-Means (Verdad Operativa)")

add_bullet_text(
    slide, Inches(0.8), Inches(1.5), Inches(4.5), Inches(3),
    [
        "K-SWEEP (K=2..6):",
        "â€¢ K=2 Ã³ptimo (Silhouette=0.232)",
        "â€¢ Estabilidad ARI=0.565",
        "",
        "CLUSTER 0 (Bajo Sedentarismo):",
        "â€¢ 402 semanas (30%)",
        "â€¢ Actividad_rel = 0.160",
        "â€¢ SuperÃ¡vit = 45.4%",
        "",
        "CLUSTER 1 (Alto Sedentarismo):",
        "â€¢ 935 semanas (70%)",
        "â€¢ Actividad_rel = 0.116",
        "â€¢ SuperÃ¡vit = 25.4%"
    ],
    font_size=16
)

# Tabla K-sweep
table_data = [
    ["K", "Silhouette", "Davies-B", "TamaÃ±os"],
    ["2", "0.232", "2.058", "{0:402, 1:935}"],
    ["3", "0.195", "1.721", "{0:685, 1:235, 2:417}"],
    ["4", "0.192", "1.422", "4 clusters"],
    ["5", "0.148", "1.444", "5 clusters"]
]

add_table(slide, Inches(5.5), Inches(1.8), 5, 4, table_data)

# ============================================================================
# SLIDES 7-10: FUNCIONES DE MEMBRESÃA (4 figuras)
# ============================================================================

mf_vars = [
    ("Actividad_relativa_p50", "Actividad Relativa (normalizada por exposiciÃ³n)"),
    ("Superavit_calorico_basal_p50", "SuperÃ¡vit CalÃ³rico Basal (% del TMB)"),
    ("HRV_SDNN_p50", "HRV SDNN (Variabilidad CardÃ­aca)"),
    ("Delta_cardiaco_p50", "Delta Cardiaco (FC_caminata - FC_reposo)")
]

for i, (var_name, var_title) in enumerate(mf_vars, start=7):
    print(f"[{i}/15] Generando slide de MF: {var_name}...")
    slide = add_content_slide(f"Funciones de MembresÃ­a: {var_title}")

    image_path = PLOTS_DIR / f"MF_{var_name}.png"
    if image_path.exists():
        add_image(slide, image_path, Inches(1), Inches(1.8), width=Inches(8))
    else:
        add_bullet_text(
            slide, Inches(2), Inches(3), Inches(6), Inches(2),
            [f"[Figura no encontrada: {image_path.name}]"],
            font_size=20
        )

# ============================================================================
# SLIDE 11: SISTEMA DIFUSO
# ============================================================================

print("[11/15] Generando slide de sistema difuso...")
slide = add_content_slide("Sistema de Inferencia Difusa (5 Reglas)")

add_bullet_text(
    slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5),
    [
        "ENTRADAS (4): Actividad_relativa, SuperÃ¡vit_calÃ³rico, HRV_SDNN, Delta_cardiaco",
        "",
        "FUNCIONES DE MEMBRESÃA: Triangulares por percentiles (p10-p25-p40, p35-p50-p65, p60-p75-p90)",
        "",
        "REGLAS MAMDANI (5 ejemplos):",
        "  R1: SI Actividad es Baja Y SuperÃ¡vit es Bajo â†’ Sedentarismo Alto",
        "  R2: SI Actividad es Alta Y SuperÃ¡vit es Alto â†’ Sedentarismo Bajo",
        "  R3: SI HRV es Baja Y Delta es Alta â†’ Sedentarismo Alto (desacondicionamiento)",
        "  R4: SI Actividad es Media Y HRV es Media â†’ Sedentarismo Medio",
        "  R5: SI Actividad es Baja Y SuperÃ¡vit es Medio â†’ Sedentarismo Medio-Alto",
        "",
        "DEFUZZIFICACIÃ“N: Centroide â†’ Score âˆˆ [0, 1]",
        "BINARIZACIÃ“N: Ï„ = 0.30 (maximiza F1-score vs. clusters)"
    ],
    font_size=17
)

# ============================================================================
# SLIDE 12: MÃ‰TRICAS GLOBALES
# ============================================================================

print("[12/15] Generando slide de mÃ©tricas globales...")
slide = add_content_slide("Resultados: MÃ©tricas de ValidaciÃ³n")

# MÃ©tricas principales (grande)
metrics_box = slide.shapes.add_textbox(
    Inches(1.5), Inches(1.8), Inches(7), Inches(2))
metrics_frame = metrics_box.text_frame
metrics_text = """F1-Score: 0.840  |  Recall: 97.6%  |  Accuracy: 74.0%
Precision: 73.7%  |  MCC: 0.294  |  Ï„ = 0.30"""
metrics_frame.text = metrics_text
metrics_frame.paragraphs[0].font.size = Pt(28)
metrics_frame.paragraphs[0].font.bold = True
metrics_frame.paragraphs[0].font.color.rgb = COLOR_ACENTO
metrics_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Tabla de matriz de confusiÃ³n
add_bullet_text(
    slide, Inches(0.8), Inches(4), Inches(4), Inches(1.5),
    [
        "MATRIZ DE CONFUSIÃ“N:",
        "  â€¢ TN = 77    (Bajo correcto)",
        "  â€¢ FP = 325  (Sobreclas. conservador)",
        "  â€¢ FN = 22    (SubclasificaciÃ³n baja)",
        "  â€¢ TP = 913  (Alto correcto)"
    ],
    font_size=16
)

# InterpretaciÃ³n
add_bullet_text(
    slide, Inches(5), Inches(4), Inches(4.2), Inches(2.5),
    [
        "INTERPRETACIÃ“N:",
        "",
        "âœ“ Alta Sensibilidad (97.6%):",
        "  Solo 22/935 semanas de Alto Sed.",
        "  pasan desapercibidas",
        "",
        "âœ“ Trade-off FP (26%):",
        "  PolÃ­tica conservadora para",
        "  screening poblacional"
    ],
    font_size=15
)

# ============================================================================
# SLIDE 13: MATRIZ DE CONFUSIÃ“N (VISUAL)
# ============================================================================

print("[13/15] Generando slide de matriz de confusiÃ³n visual...")
slide = add_content_slide("Matriz de ConfusiÃ³n (Visual)")

image_path = PLOTS_DIR / "confusion_matrix.png"
if image_path.exists():
    add_image(slide, image_path, Inches(1.5), Inches(1.8), width=Inches(7))
else:
    add_bullet_text(
        slide, Inches(2), Inches(3), Inches(6), Inches(2),
        ["[Figura no encontrada: confusion_matrix.png]"],
        font_size=20
    )

# ============================================================================
# SLIDE 14: CURVA PRECISION-RECALL Y DISTRIBUCIÃ“N
# ============================================================================

print("[14/15] Generando slide de curva PR...")
slide = add_content_slide("Curva Precision-Recall y DistribuciÃ³n de Scores")

# PR curve (izquierda)
image_path_pr = PLOTS_DIR / "pr_curve.png"
if image_path_pr.exists():
    add_image(slide, image_path_pr, Inches(
        0.5), Inches(1.8), width=Inches(4.5))

# Score distribution (derecha)
image_path_dist = PLOTS_DIR / "score_distribution_by_cluster.png"
if image_path_dist.exists():
    add_image(slide, image_path_dist, Inches(
        5.2), Inches(1.8), width=Inches(4.5))

# ============================================================================
# SLIDE 15: CONCORDANCIA POR USUARIO
# ============================================================================

print("[15/15] Generando slide de concordancia por usuario...")
slide = add_content_slide("Concordancia por Usuario")

add_bullet_text(
    slide, Inches(0.8), Inches(1.5), Inches(4.5), Inches(4),
    [
        "HETEROGENEIDAD INTER-SUJETO:",
        "",
        "Concordancia media: 70.0%",
        "Rango: 27.7% - 99.3%",
        "",
        "ALTA CONCORDANCIA (>90%):",
        "  â€¢ u1 (ale): 99.3%",
        "  â€¢ u7 (kevin): 94.7%",
        "  â†’ Patrones estables",
        "",
        "BAJA CONCORDANCIA (<50%):",
        "  â€¢ u3 (christina): 27.7%",
        "  â€¢ u8 (legarda): 44.0%",
        "  â†’ Alta variabilidad intra-semanal"
    ],
    font_size=17
)

# Tabla de concordancia
table_data = [
    ["Usuario", "Concordancia", "F1", "Recall"],
    ["u1 (ale)", "99.3%", "0.997", "1.000"],
    ["u7 (kevin)", "94.7%", "0.973", "1.000"],
    ["u6 (fidel)", "81.7%", "0.898", "0.982"],
    ["u10 (vane)", "80.9%", "0.895", "1.000"],
    ["u3 (christina)", "27.7%", "0.215", "0.875"],
    ["u8 (legarda)", "44.0%", "0.462", "0.868"]
]

add_table(slide, Inches(5.5), Inches(1.8), 7, 4, table_data)

# ============================================================================
# SLIDE 16: CONCLUSIONES
# ============================================================================

print("[16/17] Generando slide de conclusiones...")
slide = add_content_slide("Conclusiones Principales")

add_bullet_text(
    slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5),
    [
        "1. SISTEMA FUZZY VALIDADO:",
        "   Convergencia robusta con clustering K=2 (F1=0.84, Recall=97.6%)",
        "   Reglas interpretables capturan estructura real del sedentarismo",
        "",
        "2. POLÃTICA CONSERVADORA EFECTIVA:",
        "   Alta sensibilidad minimiza falsos negativos â†’ Screening poblacional",
        "   Trade-off FP aceptado (26%) con confirmaciÃ³n clÃ­nica posterior",
        "",
        "3. VARIABLES FISIOLÃ“GICAMENTE RELEVANTES:",
        "   Actividad_relativa y SuperÃ¡vit_calÃ³rico_basal (principales discriminadores)",
        "   HRV_SDNN y Delta_cardiaco (complementarios) â†’ IntegraciÃ³n multivariada",
        "",
        "4. HETEROGENEIDAD MANEJABLE:",
        "   Concordancia usuario-especÃ­fica 27.7%-99.3%",
        "   â†’ PersonalizaciÃ³n futura necesaria (Ï„ ajustable, reglas por IQR)",
        "",
        "5. TRAZABILIDAD Y REPRODUCIBILIDAD:",
        "   Pipeline documentado, auditorÃ­as de imputaciÃ³n, recalibraciÃ³n fÃ¡cil"
    ],
    font_size=16
)

# ============================================================================
# SLIDE 17: PRÃ“XIMOS PASOS
# ============================================================================

print("[17/17] Generando slide de prÃ³ximos pasos...")
slide = add_content_slide("PrÃ³ximos Pasos")

add_bullet_text(
    slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5),
    [
        "CORTO PLAZO:",
        "  â€¢ PersonalizaciÃ³n de Ï„ por usuario o subpoblaciones (sexo, rango de TMB)",
        "  â€¢ Reglas moduladas por IQR para capturar intermitencia conductual",
        "  â€¢ AnÃ¡lisis de sensibilidad de MF (variar percentiles Â±5%, medir impacto en F1)",
        "",
        "MEDIANO PLAZO:",
        "  â€¢ ValidaciÃ³n externa en nueva cohorte (â‰¥20 usuarios, â‰¥1,000 semanas)",
        "  â€¢ IntegraciÃ³n de nuevas variables: SueÃ±o (duraciÃ³n, eficiencia), estrÃ©s percibido",
        "  â€¢ Zona gris (scores 0.40-0.60) â†’ Etiqueta 'Indeterminado' + evaluaciÃ³n adicional",
        "",
        "LARGO PLAZO:",
        "  â€¢ Modelado temporal avanzado: ARIMA/LSTM para predicciÃ³n de tendencias",
        "  â€¢ ImplementaciÃ³n de dashboard clÃ­nico interactivo (FastAPI + React + Plotly)",
        "  â€¢ PublicaciÃ³n cientÃ­fica en revista de salud digital (JMIR mHealth, Digital Health)",
        "  â€¢ Despliegue clÃ­nico en programa de salud ocupacional o comunitaria"
    ],
    font_size=16
)

# ============================================================================
# SLIDE 18: AGRADECIMIENTOS
# ============================================================================

print("[18/18] Generando slide de agradecimientos...")
slide = add_title_slide(
    "Â¡Gracias por su atenciÃ³n!",
    "Preguntas y DiscusiÃ³n"
)

# Contacto
contact_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(5), Inches(9), Inches(1.5))
contact_frame = contact_box.text_frame
contact_frame.text = """Luis Ãngel MartÃ­nez
MaestrÃ­a en Ciencias, Semestre 3
luis.martinez@institution.edu

Sistema de Inferencia Difusa para EvaluaciÃ³n de Sedentarismo
Validado con F1=0.84, Recall=97.6%"""
contact_frame.paragraphs[0].font.size = Pt(16)
contact_frame.paragraphs[0].font.color.rgb = COLOR_TEXTO
contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================================
# GUARDAR PRESENTACIÃ“N
# ============================================================================

print()
print("=" * 80)
print("GUARDANDO PRESENTACIÃ“N...")
print("=" * 80)

prs.save(OUTPUT_FILE)

print()
print("=" * 80)
print("âœ… PRESENTACIÃ“N GENERADA EXITOSAMENTE")
print("=" * 80)
print()
print(f"ðŸ“ Archivo generado: {OUTPUT_FILE.name}")
print(f"ðŸ“Š Total de slides: {len(prs.slides)}")
print()
print("ðŸ“Œ CONTENIDO DE LA PRESENTACIÃ“N:")
print("   [1]  TÃ­tulo")
print("   [2]  Objetivos")
print("   [3]  Datos y Cohorte")
print("   [4]  Variables Derivadas Clave")
print("   [5]  Pipeline MetodolÃ³gico")
print("   [6]  Clustering K-Means")
print("   [7]  Funciones de MembresÃ­a: Actividad Relativa")
print("   [8]  Funciones de MembresÃ­a: SuperÃ¡vit CalÃ³rico")
print("   [9]  Funciones de MembresÃ­a: HRV SDNN")
print("   [10] Funciones de MembresÃ­a: Delta Cardiaco")
print("   [11] Sistema de Inferencia Difusa")
print("   [12] MÃ©tricas de ValidaciÃ³n")
print("   [13] Matriz de ConfusiÃ³n (Visual)")
print("   [14] Curva PR y DistribuciÃ³n de Scores")
print("   [15] Concordancia por Usuario")
print("   [16] Conclusiones Principales")
print("   [17] PrÃ³ximos Pasos")
print("   [18] Agradecimientos")
print()
print("ðŸš€ Para abrir la presentaciÃ³n, ejecuta:")
print(f"   Start-Process '{OUTPUT_FILE.name}'")
print()
print("ðŸ’¡ RECOMENDACIONES:")
print("   â€¢ Revisar cada slide antes de presentar")
print("   â€¢ Practicar timing (3-5 min por slide, total ~60 min)")
print("   â€¢ Preparar respuestas a preguntas sobre FP, heterogeneidad, Ï„ Ã³ptimo")
print()



